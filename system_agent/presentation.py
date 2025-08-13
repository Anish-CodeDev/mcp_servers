from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from gemini import extract_color_params,get_rgb
def create_slide(title,content,ppt,title_color,bg_color,content_color):
    slide_layout = ppt.slide_layouts[1] # 1-> blank slide
    slide = ppt.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    s_t = slide.shapes.title.text_frame.paragraphs[0]
    s_t.font.color.rgb = RGBColor(title_color[0],title_color[1],title_color[2])
    tf = slide.placeholders[1].text_frame

    tf.clear()
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(bg_color[0],bg_color[1],bg_color[2])

    for c in content:
        p = tf.add_paragraph()
        p.text = c
        p.font.size = Pt(30)
        p.font.color.rgb = RGBColor(content_color[0],content_color[1],content_color[2])
        p.alignment = PP_ALIGN.LEFT
#create_presentation('Applications of AI',['Defining Artificial Intelligence (AI)', 'Brief history and evolution of AI', 'The growing importance of AI in our lives', "Overview of AI\'s potential to solve complex problems"])